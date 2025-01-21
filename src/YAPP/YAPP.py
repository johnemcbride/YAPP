from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from io import BytesIO
from pptx import Presentation
import pandas as pd
from datetime import datetime, timedelta
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.dates import date2num, DateFormatter
import re


def next_thursday(today=None):
    if today is None:
        today = datetime.now()
    days_until_next_thursday = (3 - today.weekday() + 7) % 7
    if days_until_next_thursday == 0:
        days_until_next_thursday = 7
    return (today + timedelta(days=days_until_next_thursday)).date()


def busday_count_wrapper(start, end, status_datetime):
    if pd.isna(start) or pd.isna(end):
        return 0
    start_np = np.datetime64(start.date(), 'D')
    end_np = np.datetime64(min(end, status_datetime).date(), 'D')
    return np.busday_count(start_np, end_np)


def aggregate_at_level(df, level, status_datetime, hierarchy_columns):
    group_columns = hierarchy_columns[:level]

    aggregation = df.groupby(group_columns).agg(
        baseline_start=("baseline_start", "min"),
        baseline_finish=("baseline_finish", "max"),
        actual_start=("actual_start", "min"),
        actual_finish=("actual_finish", "max"),
        actual_cost=("actual_cost", "sum"),
        expected_finish=("calculated_expected_finish", "max"),
        earned_value=("earned_value", "sum"),
        planned_value=("planned_value", "sum"),
        sum_baseline_duration=("baseline_duration", "sum"),
    ).reset_index()

    for col in ["baseline_start", "baseline_finish", "actual_start", "expected_finish", "actual_finish"]:
        aggregation[col] = pd.to_datetime(aggregation[col], errors="coerce")

    def compute_baseline_duration(row):
        if pd.notna(row["baseline_start"]) and pd.notna(row["baseline_finish"]):
            return np.busday_count(np.datetime64(row["baseline_start"].date(), 'D'),
                                   np.datetime64(row["baseline_finish"].date(), 'D')) + 1
        return 0

    def compute_expected_duration(row):
        if pd.notna(row["actual_start"]) and pd.notna(row["expected_finish"]):
            return np.busday_count(np.datetime64(row["actual_start"].date(), 'D'),
                                   np.datetime64(row["expected_finish"].date(), 'D')) + 1
        return 0

    aggregation["baseline_duration"] = aggregation.apply(
        compute_baseline_duration, axis=1)
    aggregation["expected_duration"] = aggregation.apply(
        compute_expected_duration, axis=1)

    aggregation["ev_over_pv"] = np.where(
        aggregation["planned_value"] > 0,
        aggregation["earned_value"] / aggregation["planned_value"],
        0
    )

    # actual percent complete

    aggregation["actual_percent_complete"] = np.where(
        aggregation["sum_baseline_duration"] > 0,
        aggregation["earned_value"] / aggregation["sum_baseline_duration"],
        0
    )

    aggregation["expected_percent_complete"] = np.where(
        aggregation["sum_baseline_duration"] > 0,
        aggregation["planned_value"] /
        aggregation["sum_baseline_duration"],
        0
    )

    aggregation["actual_finish"] = np.where(
        aggregation["actual_percent_complete"] == 1,
        aggregation["actual_finish"],
        pd.NaT
    )

    def determine_status(row):
        if pd.isna(row["baseline_start"]) or pd.isna(row["baseline_finish"]):
            return "Unplanned"
        if pd.notna(row["actual_finish"]):
            return "Complete"
        if row["baseline_start"] <= status_datetime and pd.isna(row["actual_start"]):
            return "Late Starting"
        if row["baseline_start"] > status_datetime and pd.isna(row["actual_start"]):
            return "Not Due To Start"
        if row["ev_over_pv"] >= 0.8:
            return "In Progress"
        elif 0.3 <= row["ev_over_pv"] < 0.8:
            return "Delayed"
        else:
            return "Severely Delayed"

    aggregation["status"] = aggregation.apply(determine_status, axis=1)
    return aggregation


STATUS_COLORS = {
    "Unplanned": RGBColor(255, 255, 0),
    "In Progress": RGBColor(0, 255, 0),
    "Complete": RGBColor(0, 0, 255),
    "Delayed": RGBColor(255, 165, 0),
    "Severely Delayed": RGBColor(255, 0, 0),
    "Not Due To Start": RGBColor(169, 169, 169),
    "Late Starting": RGBColor(255, 165, 0),
}

STATUS_MPL_COLORS = {
    "Unplanned": "yellow",
    "In Progress": "green",
    "Complete": "blue",
    "Delayed": "orange",
    "Severely Delayed": "red",
    "Not Due To Start": "gray",
    "Late Starting": "orange",
}


def apply_status_colors(table, status_col_index):
    try:
        for i in range(1, len(table.rows)):
            row = table.rows[i]
            status_text = row.cells[status_col_index].text.strip(
            ) if row.cells[status_col_index].text else ""
            if status_text in STATUS_COLORS:
                row.cells[status_col_index].fill.solid()
                row.cells[status_col_index].fill.fore_color.rgb = STATUS_COLORS[status_text]
    except AttributeError as e:
        print(f"AttributeError: {e}")
    except IndexError as e:
        print(f"IndexError: {e}")


def format_date(value):
    return value.strftime("%d/%m/%Y") if pd.notna(value) else " "


def set_font_size(table, font_size):
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)


def adjust_table_height_and_font(table, row_height=0.3):
    # Set each row height to 0.3 inches
    for row in table.rows:
        row.height = Inches(row_height)
    set_font_size(table, 10)


def add_slide_title(slide, title_text):
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    else:
        textbox = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        textbox.text_frame.text = title_text


def add_free_text_box(slide):
    # Add a full-width box for free text at the top
    tx_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.8), Inches(9), Inches(0.4))
    tx_box.text_frame.text = "Status Comments:"
    return tx_box


def create_gantt_chart(gantt_data):
    gantt_data.sort_values(by="baseline_start", inplace=True)
    gantt_data = gantt_data.dropna(subset=["baseline_start", "baseline_end"])

    available_levels = [
        c for c in gantt_data.columns if re.match(r"Level_\d+", c)]
    if available_levels:
        available_levels.sort(key=lambda x: int(x.split('_')[1]))
        label_col = available_levels[-1]
    else:
        label_col = None

    fig, ax = plt.subplots(figsize=(5, 6))  # Smaller width to fit side by side
    for i, row in enumerate(gantt_data.itertuples()):
        start = row.baseline_start
        end = row.baseline_end
        if pd.notna(start) and pd.notna(end):
            dur = date2num(end) - date2num(start)
            bar_color = STATUS_MPL_COLORS.get(row.status, "black")
            ax.barh(i, dur, left=date2num(start), height=0.4,
                    color=bar_color, edgecolor='black')

    ax.set_yticks(range(len(gantt_data)))
    if label_col and label_col in gantt_data.columns:
        ylabels = gantt_data[label_col].astype(str).tolist()
    else:
        ylabels = [""] * len(gantt_data)

    ax.set_yticklabels(ylabels)
    ax.xaxis.set_major_formatter(DateFormatter('%d/%m/%Y'))
    plt.gca().invert_yaxis()
    plt.tight_layout()

    img_stream = BytesIO()
    plt.savefig(img_stream, format='png')
    plt.close(fig)
    img_stream.seek(0)
    return img_stream


def clean_slide(slide):
    # Collect all shapes and remove each one by accessing the underlying XML element
    for shape in list(slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)


def add_table_and_gantt_slide(pres, title_text, data, gantt_data, date_like_cols, status_cols):
    slide = pres.slides.add_slide(pres.slide_layouts[5])
    # Print slide layouts
    for i, master in enumerate(pres.slide_masters):
        print(f"Master {i}:")
    for j, layout in enumerate(master.slide_layouts):
        print(f"  Layout {j}: {layout.name}")
    clean_slide(slide)
    add_slide_title(slide, title_text)
    add_free_text_box(slide)

    rows, cols = data.shape

    # Table on the left
    table_left = Inches(0.5)
    table_top = Inches(1.5)
    table_width = Inches(5)
    table_height = Inches(5)

    shape = slide.shapes.add_table(
        rows+1, cols, table_left, table_top, table_width, table_height)
    table = shape.table

    # Set headers
    for col_index, col_name in enumerate(data.columns):
        table.cell(0, col_index).text = col_name

    # Reset index before iterrows to avoid issues
    data = data.reset_index(drop=True)
    for row_index, row_data in data.iterrows():
        for col_index, cell_value in enumerate(row_data):
            table.cell(row_index+1, col_index).text = str(cell_value)

    # Apply status colors if status column present
    if "Status" in data.columns:
        status_col_index = list(data.columns).index("Status")
        apply_status_colors(table, status_col_index)

    adjust_table_height_and_font(table)

    # Set column widths
    first_col_width = Inches(3)
    date_width = Inches(0.9)
    status_width = Inches(1)
    for c_i, c_name in enumerate(data.columns):
        if c_i == 0:  # first column
            table.columns[c_i].width = first_col_width
        elif c_name in date_like_cols:
            table.columns[c_i].width = date_width
        elif c_name in status_cols:
            table.columns[c_i].width = status_width
        else:
            table.columns[c_i].width = Inches(1)

    # Set each row height to 0.3 inches after population
    for r in table.rows:
        r.height = Inches(0.3)

    # Gantt chart on the right
    gantt_img = create_gantt_chart(gantt_data)
    gantt_left = Inches(9)
    gantt_top = Inches(1.5)
    gantt_width = Inches(4)
    gantt_height = Inches(5)
    slide.shapes.add_picture(gantt_img, gantt_left,
                             gantt_top, gantt_width, gantt_height)


def get_plan_data_from_csv():
    input_file = "jira_flattened_hierarchy_with_field_mapping.csv"
    return pd.read_csv(input_file)


class ComputeStatus:

    status_date = next_thursday()
    source_data_frame = pd.DataFrame()
    cleansed_data = pd.DataFrame()
    aggregates = {}

    date_columns = ["Baseline Start", "Baseline Finish",
                    "Actual Start", "Actual Finish", "Expected Finish"]

    def __init__(self, source_dataframe, status_date):
        self.source_data_frame = source_dataframe
        self.status_date = status_date

    def cleanse_input_data(self):
        def date_columns_to_datetime(df, date_columns):
            df = df.copy()
            for col in date_columns:
                df[col] = pd.to_datetime(df[col], errors="coerce")
            return df

        def rename_columns(df, date_columns):
            df = df.copy()
            for col in date_columns:
                df[col] = pd.to_datetime(df[col], errors="coerce")
            return df

        self.cleansed_data = date_columns_to_datetime(
            self.source_data_frame, self.date_columns)

    def add_calculated_fields(self):

        self.cleansed_data["baseline_start"] = self.cleansed_data["Baseline Start"]
        self.cleansed_data["baseline_finish"] = self.cleansed_data["Baseline Finish"]
        self.cleansed_data["actual_start"] = self.cleansed_data["Actual Start"]
        self.cleansed_data["actual_finish"] = self.cleansed_data["Actual Finish"]
        self.cleansed_data["expected_finish"] = self.cleansed_data["Expected Finish"]

        del self.cleansed_data["Baseline Start"]
        del self.cleansed_data["Baseline Finish"]
        del self.cleansed_data["Actual Start"]
        del self.cleansed_data["Actual Finish"]
        del self.cleansed_data["Expected Finish"]

        def busday_count_wrapper(start, end, status_datetime):
            if pd.isna(start) or pd.isna(end):
                return 0
            start_np = np.datetime64(start.date(), 'D')
            end_np = np.datetime64(min(end, status_datetime).date(), 'D')
            return np.busday_count(start_np, end_np)

        # Baseline duration
        # Working days from baseline start to finish

        def row_baseline_duration(row):
            if pd.notna(row["baseline_start"]) and pd.notna(row["baseline_finish"]):
                return np.busday_count(
                    np.datetime64(row["baseline_start"].date(), "D"),
                    np.datetime64(row["baseline_finish"].date(), "D")
                ) + 1
            return 0

        self.cleansed_data["baseline_duration"] = self.cleansed_data.apply(
            row_baseline_duration, axis=1)

        # Calculated expected finish
        # Actual finish if set, else
        # If not actually finished,
        # either the expected finish which has been inputted or the max of baseline finsih, the actual finish or status date

        self.cleansed_data['calculated_expected_finish'] = self.cleansed_data.apply(
            lambda row: (
                row['actual_finish'] if pd.notna(row['actual_finish']) else
                row['expected_finish'] if pd.notna(row['expected_finish']) else
                max(row['baseline_finish'], row['actual_finish'],
                    pd.to_datetime(self.status_date))
            ),
            axis=1)

        # Planned value
        # Essentially the time elapsed since the baseline start, up to a cap of the finish or status date

        def compute_planned_value(row, config_date):
            """
            From Excel

            =IF(Config!$B$2<$G1352,0,IF($G1352 <>"",MAX(NETWORKDAYS($G1352,MIN(Config!$B$2,$H1352)),0),0))
            =IF(status_date < baseline_start, 0, if(baseline_start not null, Max(working days between(baseline start, min(status, baseline finish))),0) )

            Replicates:

                IF config_date < start_date:
                    0
                ELSE IF start_date not blank:
                    max(NETWORKDAYS(start_date, min(config_date, end_date)), 0)
                ELSE:
                    0
            """
            start_date = row["baseline_start"]
            end_date = row["baseline_finish"]

            # If start_date is NaT/None → 0
            if pd.isna(start_date):
                return 0

            # If config_date < start_date → 0
            if pd.to_datetime(config_date) < start_date:
                return 0

            busdays = busday_count_wrapper(
                start_date, end_date, pd.to_datetime(
                    config_date)
            )+1  # as it's the start of the start date to the end of the finish date
            # Clamp to 0 if negative
            return max(busdays, 0)

        self.cleansed_data["planned_value"] = self.cleansed_data.apply(
            lambda row: compute_planned_value(row, self.status_date),
            axis=1
        )

        # Actual cost
        # How long it has taken so far or took in entirety
        # (Probably could be simplified...?)

        def compute_actual_cost(row, config_date):
            """
            From Excel

            =IF(I1352<>"",MAX(NETWORKDAYS(I1352,MIN(R1352,Config!$B$2)),0),0)
            =IF(if(actual_start not null, Max(working days between(actual start, min(calculated expected finish, status date))),0) )


            """
            start_date = row["actual_start"]
            end_date = row["calculated_expected_finish"]

            # If start_date is NaT/None → 0
            if pd.isna(start_date):
                return 0

            busdays = busday_count_wrapper(
                start_date, end_date, pd.to_datetime(
                    config_date)
            )+1  # as it's the start of the start date to the end of the finish date
            # Clamp to 0 if negative
            return max(busdays, 0)

        self.cleansed_data["actual_cost"] = self.cleansed_data.apply(
            lambda row: compute_actual_cost(row, self.status_date),
            axis=1
        )

        # Expected total cost
        # Actual start to expected finsih
        def row_expected_total_cost(row):
            if pd.notna(row["actual_start"]) and pd.notna(row["calculated_expected_finish"]):
                return np.busday_count(
                    np.datetime64(row["actual_start"].date(), "D"),
                    np.datetime64(
                        row["calculated_expected_finish"].date(), "D")
                ) + 1
            return 0
        self.cleansed_data["expected_total_cost"] = self.cleansed_data.apply(
            row_expected_total_cost, axis=1)

        # Earned value
        # if baseline start set then actual cost/expected total cost * baseline duration

        self.cleansed_data['earned_value'] = np.where(
            pd.notna(self.cleansed_data["baseline_start"]),
            self.cleansed_data["actual_cost"] /
            self.cleansed_data["expected_total_cost"] *
            self.cleansed_data['baseline_duration'],
            0
        )

        calculated_date_columns = ["baseline_start", "baseline_finish",
                                   "actual_start", "actual_finish", "calculated_expected_finish"]
        for col in calculated_date_columns:
            self.cleansed_data[col] = pd.to_datetime(
                self.cleansed_data[col], errors="coerce")

    def get_hierarchy_column_names(self):
        return [
            col for col in self.cleansed_data.columns if re.match(r"Level \d+", col)]

    def compute_aggregates(self):
        max_level = len(self.get_hierarchy_column_names())
        for level_num in range(1, max_level+1):
            self.aggregates[f"Level_{level_num}_Aggregation"] = aggregate_at_level(
                self.cleansed_data, level_num, pd.to_datetime(self.status_date), self.get_hierarchy_column_names())


def main():
    print("Calculating Status!")
    status_date = next_thursday()
    status_datetime = pd.to_datetime(status_date)
    print(f"Status Date: {status_date}")

    df = get_plan_data_from_csv()

    date_columns = ["Baseline Start", "Baseline Finish",
                    "Actual Start", "Actual Finish", "Expected Finish"]
    for col in date_columns:
        df[col] = pd.to_datetime(df[col], errors="coerce")

    df["baseline_start"] = df["Baseline Start"]
    df["baseline_end"] = df["Baseline Finish"]
    df["actual_start"] = df["Actual Start"]
    df["actual_finish"] = df["Actual Finish"]
    df["expected_finish"] = df["Expected Finish"]

    df['calculated_expected_finish'] = df.apply(
        lambda row: (
            row['actual_finish'] if pd.notna(row['actual_finish']) else
            row['expected_finish'] if pd.notna(row['expected_finish']) else
            max(row['baseline_end'], pd.to_datetime(status_date))
        ),
        axis=1
    )

    df['actual_cost'] = df.apply(
        lambda row: busday_count_wrapper(
            row['actual_start'], row['calculated_expected_finish'], status_datetime
        ),
        axis=1
    )

    # Identify hierarchy columns
    hierarchy_columns = [
        col for col in df.columns if re.match(r"Level \d+", col)]
    calculated_date_columns = ["baseline_start", "baseline_end",
                               "actual_start", "actual_finish", "calculated_expected_finish"]
    for col in calculated_date_columns:
        df[col] = pd.to_datetime(df[col], errors="coerce")

    aggregates = {}
    max_level = len(hierarchy_columns)
    for level_num in range(1, max_level+1):
        aggregates[f"Level_{level_num}_Aggregation"] = aggregate_at_level(
            df, level_num, status_datetime, hierarchy_columns)

    # -----
    output_file = "aggregated_hierarchy.xlsx"
    with pd.ExcelWriter(output_file) as writer:
        for level, agg_df in aggregates.items():
            agg_df.to_excel(writer, sheet_name=level, index=False)
    print(f"Aggregated results saved to {output_file}")

    # --------

    template_file = "Template.pptx"
    template_file = 'Cross Workstream Delivery Management 121224.pptx'
    presentation = Presentation(template_file)

    title_layout = presentation.slide_layouts[0]
    title_slide = presentation.slides.add_slide(title_layout)
    human_readable_date = status_date.strftime("%A, %d %B %Y")
    add_slide_title(title_slide, f"Status Date: {human_readable_date}")

    columns_to_drop = ["baseline_duration", "expected_duration",
                       "planned_value", "earned_value", "ev_over_pv", "actual_cost"]

    rename_map = {
        "baseline_start": "Baseline Start",
        "baseline_end": "Baseline Finish",
        "actual_start": "Actual Start",
        "actual_finish": "Actual Finish",
        "actual_cost": "Actual Cost",
        "expected_finish": "Expected Finish",
        "status": "Status"
    }

    date_like_cols = ["Baseline Start", "Baseline Finish",
                      "Actual Start", "Actual Finish", "Expected Finish"]

    status_cols = ["Status"]

    # Create a Level 2 slide if exists
    if "Level_2_Aggregation" in aggregates:
        level = 2
        data = aggregates["Level_2_Aggregation"].copy()

        for c in columns_to_drop:
            if c in data.columns:
                data.drop(columns=[c], inplace=True)

        # drop redundant Edplore part 1 column
        data = data.iloc[:, 1:]

        for c in data.columns:
            if re.match(r"Level_\d+", c):
                new_col = c.replace("_", " ")
                data.rename(columns={c: new_col}, inplace=True)
            elif c in rename_map:
                data.rename(columns={c: rename_map[c]}, inplace=True)

        for dc in date_like_cols:
            if dc in data.columns:
                data[dc] = data[dc].apply(format_date)

        gantt_data = aggregates["Level_2_Aggregation"].copy()
        add_table_and_gantt_slide(
            presentation, f"Aggregates - Level {level}", data, gantt_data, date_like_cols, status_cols)

    # Level 3 slides per unique Level_2 value
    if "Level_3_Aggregation" in aggregates:
        level = 3
        full_data = aggregates["Level_3_Aggregation"].copy()
        if "Level 2 Summary" in full_data.columns:
            unique_areas = full_data["Level 2 Summary"].unique()
            for area in unique_areas:
                subset = full_data[full_data["Level 2 Summary"] == area].copy()
                if subset.empty or subset.shape[1] == 0:
                    continue

                for c in columns_to_drop:
                    if c in subset.columns:
                        subset.drop(columns=[c], inplace=True)

                # Rename columns
                for c in subset.columns:
                    if re.match(r"Level \d+", c):
                        pass
                    elif c in rename_map:
                        subset.rename(columns={c: rename_map[c]}, inplace=True)

                subset.drop(
                    columns=["Level 1 Summary", "Level 2 Summary"], inplace=True, errors="ignore")

                for dc in date_like_cols:
                    if dc in subset.columns:
                        subset[dc] = subset[dc].apply(format_date)

                gantt_data = full_data[full_data["Level 2 Summary"] == area].copy(
                )
                add_table_and_gantt_slide(
                    presentation, f"Level 3 - {area}", subset, gantt_data, date_like_cols, status_cols)

    output_pptx = "project_aggregates_gantt.pptx"
    presentation.save(output_pptx)
    print(f"PowerPoint saved as {output_pptx}")


if __name__ == "__main__":
    main()
